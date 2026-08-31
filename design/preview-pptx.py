#!/usr/bin/env python3
"""Draw a .pptx to PNGs so a native-text deck can actually be looked at.

    python3 preview-pptx.py exports/deck.pptx [outdir]

LibreOffice cannot open a pptx in this sandbox, so there is otherwise no way
to see what a generated deck looks like. This reads the real file — the one
that ships — rather than a parallel HTML mock, so what it shows is what is in
the package.

It is a previewer, not a renderer: text wrapping is greedy and PowerPoint's
line breaking will differ slightly. Trust it for layout, collision and balance;
do not trust it to the pixel.
"""

import os
import sys
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Emu

SCALE = 150  # px per inch

FACES = {}


def _find(query):
    """Ask fontconfig rather than guessing filenames — a wrong guess falls back
    to a bitmap face and every size in the preview silently collapses."""
    import subprocess
    out = subprocess.run(['fc-match', query, '-f', '%{file}'],
                         capture_output=True, text=True).stdout.strip()
    return out or None


def face(size_px, bold):
    key = (size_px, bold)
    if key not in FACES:
        f = _find('Inter:bold' if bold else 'Inter')
        FACES[key] = ImageFont.truetype(f, size_px) if f else ImageFont.load_default()
    return FACES[key]


def rgb(color, default=(0, 0, 0)):
    try:
        v = color.rgb
        return (v[0], v[1], v[2])
    except Exception:
        return default


def wrap(draw, text, fnt, max_w):
    lines = []
    for para in text.split('\n'):
        words, cur = para.split(' '), ''
        for w in words:
            trial = (cur + ' ' + w).strip()
            if draw.textlength(trial, font=fnt) <= max_w or not cur:
                cur = trial
            else:
                lines.append(cur)
                cur = w
        lines.append(cur)
    return lines


def render(path, outdir):
    prs = Presentation(path)
    W = int(Emu(prs.slide_width).inches * SCALE)
    H = int(Emu(prs.slide_height).inches * SCALE)
    os.makedirs(outdir, exist_ok=True)
    made = []

    for idx, slide in enumerate(prs.slides, 1):
        try:
            bg = rgb(slide.background.fill.fore_color, (245, 245, 247))
        except Exception:
            bg = (245, 245, 247)
        im = Image.new('RGB', (W, H), bg)
        d = ImageDraw.Draw(im)

        for sh in slide.shapes:
            x = Emu(sh.left).inches * SCALE
            y = Emu(sh.top).inches * SCALE
            w = Emu(sh.width).inches * SCALE
            h = Emu(sh.height).inches * SCALE
            kind = str(sh.shape_type)

            if kind.startswith('PICTURE'):
                try:
                    from io import BytesIO
                    pic = Image.open(BytesIO(sh.image.blob)).convert('RGBA')
                    pic = pic.resize((max(1, int(w)), max(1, int(h))))
                    im.paste(pic, (int(x), int(y)), pic)
                except Exception:
                    d.rectangle([x, y, x + w, y + h], outline=(200, 200, 200))
                continue

            if kind.startswith('AUTO_SHAPE') or kind.startswith('FREEFORM'):
                try:
                    fill = rgb(sh.fill.fore_color, None) if sh.fill.type is not None else None
                except Exception:
                    fill = None
                if fill:
                    if 'OVAL' in sh.name.upper() or (abs(w - h) < 2 and w < 80):
                        d.ellipse([x, y, x + w, y + h], fill=fill)
                    else:
                        d.rounded_rectangle([x, y, x + w, y + h], radius=min(26, w / 6), fill=fill)

            # pptxgenjs gives every shape a text frame, so a zero-height rule
            # cannot be told from a text box by that alone — check the height.
            blank = not sh.has_text_frame or not sh.text_frame.text.strip()
            if kind.startswith('LINE') or (h < 2 and blank):
                try:
                    lc = rgb(sh.line.color, (201, 201, 206))
                except Exception:
                    lc = (201, 201, 206)
                d.line([x, y, x + w, y + h], fill=lc, width=2)

            if not sh.has_text_frame or not sh.text_frame.text.strip():
                continue

            cy = y
            block = []
            for para in sh.text_frame.paragraphs:
                runs = [r for r in para.runs if r.text]
                if not runs:
                    continue
                txt = ''.join(r.text for r in runs)
                r0 = runs[0]
                size = int((r0.font.size.pt if r0.font.size else 18) * SCALE / 72)
                fnt = face(size, bool(r0.font.bold))
                col = rgb(r0.font.color, (29, 29, 31)) if r0.font.color and r0.font.color.type is not None else (29, 29, 31)
                bullet = '\u00b7   ' if '<a:buChar' in para._p.xml else ''
                # Use the paragraph's real leading and space-after. Guessing a
                # multiple of the point size under-reports the height and the
                # preview then shows text fitting that PowerPoint overflows.
                lead = (para.line_spacing.pt if hasattr(para.line_spacing, 'pt')
                        else (para.line_spacing * (r0.font.size.pt if r0.font.size else 18)
                              if para.line_spacing else (r0.font.size.pt if r0.font.size else 18) * 1.2))
                lead_px = lead * SCALE / 72
                after_px = (para.space_after.pt * SCALE / 72) if para.space_after else 0
                lines = wrap(d, bullet + txt, fnt, w)
                for j, ln in enumerate(lines):
                    block.append((ln, fnt, col, lead_px + (after_px if j == len(lines) - 1 else 0),
                                  para.alignment))

            total = sum(adv for _, _, _, adv, _ in block)
            anchor = str(sh.text_frame.vertical_anchor or '')
            if 'MIDDLE' in anchor:
                cy = y + (h - total) / 2

            for ln, fnt, col, adv, align in block:
                tw = d.textlength(ln, font=fnt)
                tx = x
                if align is not None and 'CENTER' in str(align):
                    tx = x + (w - tw) / 2
                d.text((tx, cy), ln, font=fnt, fill=col)
                cy += adv

        out = os.path.join(outdir, 'slide-%02d.png' % idx)
        im.save(out)
        made.append(out)

    print('rendered %d slides -> %s' % (len(made), outdir))
    return made


if __name__ == '__main__':
    render(sys.argv[1], sys.argv[2] if len(sys.argv) > 2 else 'exports/preview')
